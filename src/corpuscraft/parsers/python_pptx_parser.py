from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.table import Table

from corpuscraft.config import ParserConfig
from corpuscraft.models import ParsedDocument
from corpuscraft.parsers.base import BaseParser

logger = logging.getLogger(__name__)

_TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}


def _is_title(shape: BaseShape) -> bool:
    if not shape.is_placeholder:
        return False
    try:
        return shape.placeholder_format.type in _TITLE_PLACEHOLDERS
    except (AttributeError, ValueError):
        return False


def _text_frame_to_markdown(shape: BaseShape) -> str:
    if not shape.has_text_frame:
        return ""

    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        # level 0 paragraphs are top-level; deeper levels become nested bullets
        if para.level > 0:
            indent = "  " * (para.level - 1)
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


def _table_to_markdown(table: Table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]

    lines = ["| " + " | ".join(rows[0]) + " |"]
    lines.append("| " + " | ".join(["---"] * width) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _slide_to_markdown(slide: Slide, idx: int) -> tuple[str, dict]:
    title_text = ""
    body_parts: list[str] = []
    n_tables = 0
    n_pictures = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n_pictures += 1
            continue

        if shape.has_table:
            md = _table_to_markdown(shape.table)
            if md:
                body_parts.append(md)
                n_tables += 1
            continue

        if shape.has_text_frame:
            md = _text_frame_to_markdown(shape)
            if not md:
                continue
            if _is_title(shape) and not title_text:
                title_text = md.splitlines()[0]
            else:
                body_parts.append(md)

    header = f"## Slide {idx}"
    if title_text:
        header = f"## Slide {idx}: {title_text}"

    parts = [header]
    parts.extend(body_parts)

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            quoted = "\n".join(f"> {line}" for line in notes.splitlines())
            parts.append(f"<!-- speaker notes -->\n{quoted}")

    return "\n\n".join(parts), {
        "tables": n_tables,
        "pictures": n_pictures,
        "has_notes": slide.has_notes_slide,
    }


class PythonPptxParser(BaseParser):
    """Lightweight PPTX parser using python-pptx.

    Each slide becomes a `## Slide N: Title` section, with body text,
    bullet hierarchy, tables, and speaker notes preserved as markdown.
    """

    def __init__(self, config: ParserConfig):
        self.config = config

    def parse_file(self, path: Path) -> ParsedDocument:
        logger.info(f"Parsing {path.name} with python-pptx")

        prs = Presentation(str(path))
        slide_md: list[str] = []
        total_tables = 0
        total_pictures = 0
        slides_with_notes = 0

        for i, slide in enumerate(prs.slides, start=1):
            md, stats = _slide_to_markdown(slide, i)
            slide_md.append(md)
            total_tables += stats["tables"]
            total_pictures += stats["pictures"]
            if stats["has_notes"]:
                slides_with_notes += 1

        content = "\n\n".join(slide_md)

        core = prs.core_properties
        metadata = {
            "title": core.title or "",
            "author": core.author or "",
            "created": core.created.isoformat() if core.created else "",
            "modified": core.modified.isoformat() if core.modified else "",
            "slide_count": len(prs.slides),
            "table_count": total_tables,
            "picture_count": total_pictures,
            "slides_with_notes": slides_with_notes,
            "file_size": path.stat().st_size,
        }

        return ParsedDocument(
            content=content,
            source_path=path,
            pipeline="python_pptx",
            metadata=metadata,
        )
